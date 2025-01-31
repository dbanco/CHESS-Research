#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
spotfetch.visualization

Tools for visualizing X-ray diffraction data and results.

Created on: Fri Nov 15 23:00:28 2024
Author: Daniel Banco
"""
import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import glob
import spotfetch as sf
import spotfetch.detectors as dt
from pptx import Presentation
from pptx.util import Inches

def roiAdjacent(ind,tth,eta,frame,omFrms,timeFrms,params,dataDir):
    
    fig, axes = plt.subplots(len(omFrms), len(timeFrms), figsize=(len(omFrms), len(timeFrms)))
    
    # add enumerate HERR ND THEN coordinate subplots
    for i,frm in enumerate(omFrms):
        frm = sf.wrapFrame(frm)
        
        for j,t in enumerate(timeFrms):
            fnames = sf.timeToFile(t,dataDir)
            roi = dt.loadPolarROI(fnames,tth,eta,frm,params)
            img = axes[i,j].imshow(roi)
            fig.colorbar(img, ax=axes[i,j])
            axes[i,j].set_title(f'$\omega={frm}$, t={t}')

def plotROIs(roi_list,num_cols = 5):
    # Create subplots
    num_images = len(roi_list)
    num_rows = (num_images + num_cols - 1) // num_cols  # Calculate number of rows needed

    fig, axes = plt.subplots(num_rows, num_cols, figsize=(10, 10))

    # Plot each image in a subplot
    for i, ax in enumerate(axes.flat):
        if i < num_images:
            ax.imshow(roi_list[i])
            ax.axis('off')  # Turn off axis
            ax.set_title(f'Spot{i+1}')  # Set title for each subplot
            # Add colorbars
            pcm = ax.pcolormesh(roi_list[i])
            fig.colorbar(pcm,ax=ax)
                    
    # Adjust layout to prevent overlapping
    plt.tight_layout()
    
    plt.show()
    
    return fig

def plotSpotWedges(spotData,fnames,frame_i,params,grains=[],detectFrame=[]):
    tths = spotData['tths']
    etas = spotData['etas']     
    frms = spotData['ome_idxs']
    frms_i = frms - 2 
    
    # Get spot indices at frame
    if grains == []:
        spotInds = np.where(frms_i == frame_i)[0]
    else:
        spotInds = sf.findSpots(spotData,grains=grains,frm=frame_i+2)
    
    roiSize = params['roiSize']
    imSize = params['imSize']
    center = (imSize[0]//2,imSize[1]//2,)
        
    b = sf.loadImg(fnames,params,frame_i)
    
    fig = plt.figure()
    # b[b>100] = 100
    plt.imshow(b)
    plt.clim(np.median(b),np.median(b)+20)
    plt.colorbar()
    
    for ind in spotInds:
        # 0. Load spot information
        tth = tths[ind]
        eta = etas[ind]
        print(eta)
        print(tth)
        # 1. Load YAML data
        detectDist, mmPerPixel, ff_trans, ff_tilt = dt.loadYamlData(params,tth,eta)
        
        r = detectDist*np.tan(tth)/mmPerPixel
        x = -r*np.cos(eta)
        y = r*np.sin(eta)

        [x,y] = sf.applyTilt(x,y,ff_tilt,detectDist)

        eta = np.arctan2(y,x)
        
        r = np.sqrt(x**2 + y**2)
        tth = np.arctan(r*mmPerPixel/detectDist)
        print(eta)
        print(tth)

        # 2. Construct rad, eta domain
        rad_dom, eta_dom = dt.polarDomain(detectDist, mmPerPixel, tth, eta, roiSize)
   
        rad = np.round(detectDist*np.tan(tth)/mmPerPixel)
        
        wedge = patches.Wedge([center[1],center[0]],rad+roiSize[0]/2,\
                180/np.pi*eta_dom[0],180/np.pi*eta_dom[-1],\
                linewidth=1,width=roiSize[1],fill=0,color='r')
            
        plt.gca().add_patch(wedge)
        # x = rad*np.cos(eta) + center[1]
        # y = rad*np.sin(eta) + center[0]
        # plt.plot(x,y,color='r',marker='x')
        
    return fig

def scatterOmeg(trackData):
    """
    Visualizes omega values from track data in scatter plots.
    
    This function generates a series of scatter plots where each subplot corresponds 
    to a specific spot (indexed by `k`) and visualizes the omega values extracted 
    from frame data for each time step `t` in `trackData`.
    
    Parameters:
    -----------
    trackData : list of lists of lists of dict
        A nested list where:
        - Outer list represents time steps (`t`).
        - Each element is a list containing `K` sublists, each corresponding to a spot.
        - Each sublist contains dictionaries with a `'frm'` key, 
          which is used to calculate the omega value.
    
    Returns:
    --------
    matplotlib.figure.Figure
        A matplotlib Figure object containing the scatter plots.
    
    Notes:
    ------
    - The function assumes that `sf.frameToOmega` is defined and converts frame 
      information to omega values.
    - Only non-empty sublists are processed.
    - The x-axis represents the time step (`t`), and the y-axis represents the omega value.
    - The figure contains one subplot per spot, with a maximum x-axis limit of 20.
    
    Example:
    --------
    >>> fig = scatterOmeg(trackData)
    >>> fig.show()
    """
    T = len(trackData)
    K = len(trackData[0])
    
    fig, axes = plt.subplots(K, 1, figsize=(K, 1))
    fig.subplots_adjust(hspace=0.5) 
    
    for t in range(T):
        if trackData[t] != []:
            for k in range(K):
                L = len(trackData[t][k])
                if L > 0:
                    for j in range(L):
                        omega = sf.frameToOmega(trackData[t][k][j]['frm'])
                        axes[k].scatter(t,omega,marker='s',color='b')
                        axes[k].set_title(f'spot {k}')
                        axes[k].set_xlim((0,20))

    return fig

def roiTrackVisual(spotInd,spotData,dome,num_cols,scanRange,dataPath,trackPath,truthPath,spotFiles,params):

    T = len(scanRange)
    N_ome = 2*dome+1
    
    num_figs = int(np.ceil(T/num_cols))
    
    fig_list = []
    axes_list = []
    
    # Create a figure and a set of subplots
    for i in range(num_figs):
        fig, axes = plt.subplots(N_ome, num_cols, figsize=(20, 15))
        
        # Remove x and y ticks for clarity
        for ax_row in axes:
            for ax in ax_row:
                ax.set_xticks([])
                ax.set_yticks([])
        
        fig_list.append(fig)
        axes_list.append(axes)
        
        i1 = 0 + i*num_cols
        i2 = num_cols-1 + i*num_cols
        j1 = scanRange[i1]
        if i2 >= T:
            j2 = scanRange[-1]
        else:
            j2 = scanRange[i2]
        # Add common labels
        fig_list[i].text(0.04, 0.5, r'$\omega$ frame', va='center', rotation='vertical', fontsize=24)
        fig_list[i].text(0.5, 0.04, 'Scan #', ha='center', fontsize=24)
        fig_list[i].text(0.5, 0.95, f'Spot {spotInd}, Scans {j1}-{j2}', ha='center', fontsize=32)      
        fig_list[i].subplots_adjust(wspace=0.05, hspace=0.01)
    
    # Load track data
    track_file = os.path.join(trackPath,f'trackData_{spotInd}.pkl')
    if os.path.exists(track_file):
        with open(track_file, 'rb') as f:
            outData = pickle.load(f)
            track_data = outData['trackData']
            print('Track loaded')
    else:
        track_data = []
        trackFound = False
        
    # Load truth data 
    truth_file = os.path.join(trackPath,f'truthData_{spotInd}.pkl')
    if os.path.exists(truth_file):
        with open(truth_file, 'rb') as f:
            truth_data = pickle.load(f)
            print('Truth loaded')
    else:   
        T = len(scanRange)
        truth_data = []
        truthFound = False
        
    # Load sim truth
    spotDataList = []
    for i in range(len(spotFiles)):
        if os.path.exists(spotFiles[i]):
            spotDataList.append(np.load(spotFiles[i]))
    
    # Get initial spot location
    x = spotData['Xm'][spotInd]
    y = spotData['Ym'][spotInd]
    eta0, tth0 = sf.xyToEtaTthRecenter(x,y,params)
    frm0 = spotData['ome_idxs'][spotInd]
    print(spotData['omes'][spotInd])
    
    etaRoi = eta0
    tthRoi = tth0
    frmRange = np.arange(frm0-dome,frm0+dome+1)
    spot_id = sf.getSpotID(spotData,spotInd)
    
    for scan_ind in range(T):
        for om_ind in range(N_ome):
            i = int(np.floor(scan_ind/num_cols))
            j = np.mod(scan_ind,num_cols)
            ax = axes_list[i][om_ind,j]
            scan = scanRange[scan_ind]
            frm = sf.wrapFrame(frmRange[om_ind])
            
            print(f'Scan {scan}, Frame {frm}')
            # 1. Show ROI
            sf.showROI(ax,dataPath, scan, frm, tthRoi, etaRoi, params)
            # sf.showInitial(ax,etaRoi,tthRoi,eta0,tth0,params)
            
            # 2. Show the track, truth, sim_truth if they exist
            if len(track_data) > 0:
                [trackFound, om_ind1] = sf.checkTrack(track_data,scan_ind,scan,frm)
            if len(truth_data) > 0:
                [truthFound, om_ind2] = sf.checkTruth(truth_data,scan_ind,scan,frm)  
            if len(spotDataList) > 0:
                m_ind = sf.matchSpotID(spotDataList[scan_ind],spot_id,spotInd)
                if not np.isnan(m_ind):
                    if frm == spotDataList[scan_ind]['ome_idxs'][m_ind]:
                        # Plot the sim truth
                        x = spotDataList[scan_ind]['Xm'][m_ind]
                        y = spotDataList[scan_ind]['Ym'][m_ind]
                        etaTrue, tthTrue = sf.xyToEtaTthRecenter(x,y,params)
                        [y1, x1] = sf.etaTthToPix(etaTrue,tthTrue,etaRoi,tthRoi,params)
                        if (y1 >= 0) and (y1 < params['roiSize'][0]) and\
                           (x1 >= 0) and (x1 < params['roiSize'][1]):
                            ax.plot(x1,y1,marker='o',fillstyle='none',markersize=16,color='cyan')
                            ax.plot(x1,y1,marker='s',fillstyle='none',markersize=36,color='cyan')
                    
            # Add tracks and truth if available
            if trackFound:
                track = track_data[scan_ind][om_ind1].copy()
                sf.showTrack(ax,track,etaRoi,tthRoi,\
                          eta0,tth0,params)
            if truthFound:
                truth = truth_data[scan_ind][om_ind2].copy()
                sf.showTruth(ax,truth,etaRoi,tthRoi,params)
            
            # Label plots
            if om_ind == N_ome-1:
                ax.set_xlabel(f'{scan}')
            if scan_ind == 0:
                ax.set_ylabel(f'{frm}')
    return fig_list

def initRoiFigs(scanRange,num_cols,N_ome,spotInd):
    T = len(scanRange)
    num_figs = int(np.ceil(T/num_cols))
    
    fig_list = []
    axes_list = []
    
    # Create a figure and a set of subplots
    for i in range(num_figs):
        fig, axes = plt.subplots(N_ome, num_cols, figsize=(20, 15))
        
        # Remove x and y ticks for clarity
        for ax_row in axes:
            for ax in ax_row:
                ax.set_xticks([])
                ax.set_yticks([])
        
        fig_list.append(fig)
        axes_list.append(axes)
        
        i1 = 0 + i*num_cols
        i2 = num_cols-1 + i*num_cols
        j1 = scanRange[i1]
        if i2 >= T:
            j2 = scanRange[-1]
        else:
            j2 = scanRange[i2]
        # Add common labels
        fig_list[i].text(0.04, 0.5, r'$\omega$ frame', va='center', rotation='vertical', fontsize=24)
        fig_list[i].text(0.5, 0.04, 'Scan #', ha='center', fontsize=24)
        fig_list[i].text(0.5, 0.95, f'Spot {spotInd}, Scans {j1}-{j2}', ha='center', fontsize=32)      
        fig_list[i].subplots_adjust(wspace=0.05, hspace=0.01)
        
        return fig_list, axes_list

def showRoi3D(spotInd,spotData,num_cols,scanRange,dataFileSequence,trackPath,truthPath,spotFiles,params):
    
    # Load track data
    track_data = sf.loadTrackData(trackPath,spotInd)
        
    # Load truth data 
    truth_data = sf.loadTruthData(trackPath,spotInd)
    
    # Load sim truth
    spotDataList = []
    for i in range(len(spotFiles)):
        if os.path.exists(spotFiles[i]):
            spotDataList.append(np.load(spotFiles[i]))
    
    # Get initial spot location
    x = spotData['Xm'][spotInd]
    y = spotData['Ym'][spotInd]
    eta0, tth0 = sf.xyToEtaTthRecenter(x,y,params)
    frm0 = spotData['ome_idxs'][spotInd]
    print(spotData['omes'][spotInd])
    
    spot_id = sf.getSpotID(spotData,spotInd)
    
    # Initialize figures
    N_ome = params['roiSize'][2]
    dome = (N_ome-1)/2
    frmRange = np.arange(frm0-dome,frm0+dome+1)

    fig_list, axes_list = sf.initRoiFigs(scanRange,num_cols,N_ome,spotInd)
    
    for scan_ind in range(len(scanRange)):
        # Load roi data
        scan = scanRange[scan_ind]
        roi3D = sf.loadPolarROI3D(dataFileSequence[scan_ind], tth0, eta0, frm0, params)
        
        trackFound = False
        if (track_data[scan_ind] != False):
           trackFound = True
           track = track_data[scan_ind].copy()
           etaTrack = track['eta']
           tthTrack = track['tth']
           frmTrack = round(track['frm'])
           y_track, x_track = sf.etaTthToPix(etaTrack, tthTrack, eta0, tth0, params)
           inds = np.where(np.sum(track['blob'],axis=(0,1)))
           frm1 = sf.indToFrame(np.min(inds),frmTrack,params['roiSize'][2])
           frm2 = sf.indToFrame(np.max(inds),frmTrack,params['roiSize'][2])
           
        for om_ind in range(N_ome):
            i = int(np.floor(scan_ind/num_cols))
            j = np.mod(scan_ind,num_cols)
            ax = axes_list[i][om_ind,j]
            frm = sf.wrapFrame(frmRange[om_ind])
            
            print(f'Scan {scan}, Frame {frm}')
            # 1. Show ROI
            sf.plotROI(ax,roi3D[:,:,om_ind],tth0, eta0, params)
            # sf.showInitial(ax,etaRoi,tthRoi,eta0,tth0,params)
            
            # 2. Show the track, truth, sim_truth if they exist
            if trackFound:
                if frm == frm1:
                    # Plot track square 1
                    ax.plot(x_track,y_track,marker='s',fillstyle='none',markersize=36,color='red')
                elif frm == frm2:
                    # Plot track square 2
                    ax.plot(x_track,y_track,marker='s',fillstyle='none',markersize=36,color='red')
                if (frm >= frm1) and (frm <= frm2):
                    # Plot track circle
                    ax.plot(x_track,y_track,marker='o',fillstyle='none',markersize=16,color='red')
                
            if len(truth_data) > 0:
                [truthFound, om_ind2] = sf.checkTruth(truth_data,scan_ind,scan,frm)  
            if len(spotDataList) > 0:
                m_ind = sf.matchSpotID(spotDataList[scan_ind],spot_id,spotInd)
                if not np.isnan(m_ind):
                    if frm == spotDataList[scan_ind]['ome_idxs'][m_ind]:
                        # Plot the sim truth
                        x = spotDataList[scan_ind]['Xm'][m_ind]
                        y = spotDataList[scan_ind]['Ym'][m_ind]
                        etaTrue, tthTrue = sf.xyToEtaTthRecenter(x,y,params)
                        [y1, x1] = sf.etaTthToPix(etaTrue,tthTrue,eta0,tth0,params)
                        if (y1 >= 0) and (y1 < params['roiSize'][0]) and\
                           (x1 >= 0) and (x1 < params['roiSize'][1]):
                            ax.plot(x1,y1,marker='o',fillstyle='none',markersize=16,color='cyan')
                            ax.plot(x1,y1,marker='s',fillstyle='none',markersize=36,color='cyan')
                            
            if len(truth_data) > 0:
                truth = truth_data[scan_ind][om_ind2].copy()
                sf.showTruth(ax,truth,eta0,tth0,params)
            
            # Label plots
            if om_ind == N_ome-1:
                ax.set_xlabel(f'{scan}')
            if scan_ind == 0:
                ax.set_ylabel(f'{frm}')
    return fig_list

def showRoiBlobs(spotInd,spotData,num_cols,scanRange,dataFileSequence,trackPath,truthPath,spotFiles,params):
    
    # Load track data
    track_data = sf.loadTrackData(trackPath,spotInd)
        
    # Load truth data 
    truth_data = sf.loadTruthData(trackPath,spotInd)
    
    # Load sim truth
    spotDataList = []
    for i in range(len(spotFiles)):
        if os.path.exists(spotFiles[i]):
            spotDataList.append(np.load(spotFiles[i]))
    
    # Get initial spot location
    x = spotData['Xm'][spotInd]
    y = spotData['Ym'][spotInd]
    eta0, tth0 = sf.xyToEtaTthRecenter(x,y,params)
    frm0 = spotData['ome_idxs'][spotInd]
    print(spotData['omes'][spotInd])
    
    spot_id = sf.getSpotID(spotData,spotInd)
    
    # Initialize figures
    N_ome = params['roiSize'][2]
    dome = (N_ome-1)/2
    frmRange = np.arange(frm0-dome,frm0+dome+1)

    fig_list, axes_list = sf.initRoiFigs(scanRange,num_cols,N_ome,spotInd)
    
    for scan_ind in range(len(scanRange)):
        # Load roi data
        scan = scanRange[scan_ind]
        roi3D = sf.loadPolarROI3D(dataFileSequence[scan_ind], tth0, eta0, frm0, params)
        blobs,_,_ = sf.detectBlobDoG(roi3D,params)
        
        trackFound = False
        if (track_data[scan_ind] != False):
           trackFound = True
           track = track_data[scan_ind].copy()
           etaTrack = track['eta']
           tthTrack = track['tth']
           frmTrack = round(track['frm'])
           y_track, x_track = sf.etaTthToPix(etaTrack, tthTrack, eta0, tth0, params)
           inds = np.where(np.sum(track['blob'],axis=(0,1)))
           frm1 = sf.indToFrame(np.min(inds),frmTrack,params['roiSize'][2])
           frm2 = sf.indToFrame(np.max(inds),frmTrack,params['roiSize'][2])
           
        for om_ind in range(N_ome):
            i = int(np.floor(scan_ind/num_cols))
            j = np.mod(scan_ind,num_cols)
            ax = axes_list[i][om_ind,j]
            frm = sf.wrapFrame(frmRange[om_ind])
            
            print(f'Scan {scan}, Frame {frm}')
            # 1. Show ROI
            sf.plotROI(ax,blobs[:,:,om_ind],tth0, eta0, params)
            # roiPlt.set_clim(0, 10)
            # sf.showInitial(ax,etaRoi,tthRoi,eta0,tth0,params)
            
            # 2. Show the track, truth, sim_truth if they exist
            if trackFound:
                if frm == frm1:
                    # Plot track square 1
                    ax.plot(x_track,y_track,marker='s',fillstyle='none',markersize=36,color='red')
                elif frm == frm2:
                    # Plot track square 2
                    ax.plot(x_track,y_track,marker='s',fillstyle='none',markersize=36,color='red')
                if (frm >= frm1) and (frm <= frm2):
                    # Plot track circle
                    ax.plot(x_track,y_track,marker='o',fillstyle='none',markersize=16,color='red')
                
            if len(truth_data) > 0:
                [truthFound, om_ind2] = sf.checkTruth(truth_data,scan_ind,scan,frm)  
            if len(spotDataList) > 0:
                m_ind = sf.matchSpotID(spotDataList[scan_ind],spot_id,spotInd)
                if not np.isnan(m_ind):
                    if frm == spotDataList[scan_ind]['ome_idxs'][m_ind]:
                        # Plot the sim truth
                        x = spotDataList[scan_ind]['Xm'][m_ind]
                        y = spotDataList[scan_ind]['Ym'][m_ind]
                        etaTrue, tthTrue = sf.xyToEtaTthRecenter(x,y,params)
                        [y1, x1] = sf.etaTthToPix(etaTrue,tthTrue,eta0,tth0,params)
                        if (y1 >= 0) and (y1 < params['roiSize'][0]) and\
                           (x1 >= 0) and (x1 < params['roiSize'][1]):
                            ax.plot(x1,y1,marker='o',fillstyle='none',markersize=16,color='cyan')
                            ax.plot(x1,y1,marker='s',fillstyle='none',markersize=36,color='cyan')
                            
            if len(truth_data) > 0:
                truth = truth_data[scan_ind][om_ind2].copy()
                sf.showTruth(ax,truth,eta0,tth0,params)
            
            # Label plots
            if om_ind == N_ome-1:
                ax.set_xlabel(f'{scan}')
            if scan_ind == 0:
                ax.set_ylabel(f'{frm}')
    return fig_list

def showRoiTrackBlobs(spotInd,spotData,num_cols,scanRange,dataFileSequence,trackPath,truthPath,spotFiles,params):
    
    # Load track data
    track_data = sf.loadTrackData(trackPath,spotInd)
        
    # Load truth data 
    truth_data = sf.loadTruthData(trackPath,spotInd)
    
    # Load sim truth
    spotDataList = []
    for i in range(len(spotFiles)):
        if os.path.exists(spotFiles[i]):
            spotDataList.append(np.load(spotFiles[i]))
    
    # Get initial spot location
    x = spotData['Xm'][spotInd]
    y = spotData['Ym'][spotInd]
    frm0 = spotData['ome_idxs'][spotInd]
    print(spotData['omes'][spotInd])
    
    spot_id = sf.getSpotID(spotData,spotInd)
    
    # Initialize figures
    N_ome = params['roiSize'][2]
    dome = (N_ome-1)/2
    frmRange = np.arange(frm0-dome,frm0+dome+1)

    fig_list, axes_list = sf.initRoiFigs(scanRange,num_cols,N_ome,spotInd)
    
    for scan_ind in range(len(scanRange)):
        # Load roi data
        scan = scanRange[scan_ind]        
        trackFound = False
        if (track_data[scan_ind] != False):
           trackFound = True
           track = track_data[scan_ind].copy()
           blob = track['blob']
           etaRoi = track['etaRoi']
           tthRoi = track['tthRoi']
           etaTrack = track['eta']
           tthTrack = track['tth']
           frmTrack = round(track['frm'])
           y_track, x_track = sf.etaTthToPix(etaTrack, tthTrack, etaRoi, tthRoi, params)
           inds = np.where(np.sum(track['blob'],axis=(0,1)))
           frm1 = sf.indToFrame(np.min(inds),frmTrack,params['roiSize'][2])
           frm2 = sf.indToFrame(np.max(inds),frmTrack,params['roiSize'][2])
           
        for om_ind in range(N_ome):
            i = int(np.floor(scan_ind/num_cols))
            j = np.mod(scan_ind,num_cols)
            ax = axes_list[i][om_ind,j]
            frm = sf.wrapFrame(frmRange[om_ind])
            
            print(f'Scan {scan}, Frame {frm}')
            # 1. Show ROI
            sf.plotROI(ax,blob[:,:,om_ind],tthRoi, etaRoi, params)
            # roiPlt.set_clim(0, 10)
            # sf.showInitial(ax,etaRoi,tthRoi,eta0,tth0,params)
            
            # 2. Show the track, truth, sim_truth if they exist
            if trackFound:
                if frm == frm1:
                    # Plot track square 1
                    ax.plot(x_track,y_track,marker='s',fillstyle='none',markersize=36,color='red')
                elif frm == frm2:
                    # Plot track square 2
                    ax.plot(x_track,y_track,marker='s',fillstyle='none',markersize=36,color='red')
                if (frm >= frm1) and (frm <= frm2):
                    # Plot track circle
                    ax.plot(x_track,y_track,marker='o',fillstyle='none',markersize=16,color='red')
                
            if len(truth_data) > 0:
                [truthFound, om_ind2] = sf.checkTruth(truth_data,scan_ind,scan,frm)  
            if len(spotDataList) > 0:
                m_ind = sf.matchSpotID(spotDataList[scan_ind],spot_id,spotInd)
                if not np.isnan(m_ind):
                    if frm == spotDataList[scan_ind]['ome_idxs'][m_ind]:
                        # Plot the sim truth
                        x = spotDataList[scan_ind]['Xm'][m_ind]
                        y = spotDataList[scan_ind]['Ym'][m_ind]
                        etaTrue, tthTrue = sf.xyToEtaTthRecenter(x,y,params)
                        [y1, x1] = sf.etaTthToPix(etaTrue,tthTrue,etaRoi,tthRoi,params)
                        if (y1 >= 0) and (y1 < params['roiSize'][0]) and\
                           (x1 >= 0) and (x1 < params['roiSize'][1]):
                            ax.plot(x1,y1,marker='o',fillstyle='none',markersize=16,color='cyan')
                            ax.plot(x1,y1,marker='s',fillstyle='none',markersize=36,color='cyan')
                            
            if len(truth_data) > 0:
                truth = truth_data[scan_ind][om_ind2].copy()
                sf.showTruth(ax,truth,etaRoi,tthRoi,params)
            
            # Label plots
            if om_ind == N_ome-1:
                ax.set_xlabel(f'{scan}')
            if scan_ind == 0:
                ax.set_ylabel(f'{frm}')
    return fig_list

def roiTensor(spotInd,spotData,dome,scanRange,dataPath,trackPath,truthPath,spotFiles,params):

    T = len(scanRange)
    M_ome = 2*dome+1
    
    # Load track data
    track_file = os.path.join(trackPath,f'trackData_{spotInd}.pkl')
    if os.path.exists(track_file):
        with open(track_file, 'rb') as f:
            outData = pickle.load(f)
            track_data = outData['trackData']
            print('Track loaded')
    else:
        track_data = []
        trackFound = False
        
    # Load truth data 
    truth_file = os.path.join(trackPath,f'truthData_{spotInd}.pkl')
    if os.path.exists(truth_file):
        with open(truth_file, 'rb') as f:
            truth_data = pickle.load(f)
            print('Truth loaded')
    else:   
        T = len(scanRange)
        truth_data = []
        truthFound = False
        
    # Load sim truth
    spotDataList = []
    for i in range(len(spotFiles)):
        if os.path.exists(spotFiles[i]):
            spotDataList.append(np.load(spotFiles[i]))
    
    # Get initial spot location
    x = spotData['Xm'][spotInd]
    y = spotData['Ym'][spotInd]
    eta0, tth0 = sf.xyToEtaTthRecenter(x,y,params)
    frm0 = spotData['ome_idxs'][spotInd]
    print(spotData['omes'][spotInd])
    
    etaRoi = eta0
    tthRoi = tth0
    frmRange = np.arange(frm0-dome,frm0+dome+1)

    M_tth = params['roiSize'][0]
    M_eta = params['roiSize'][1]
    roiTensor = np.zeros((M_tth,M_eta,M_ome,T))
    
    for scan_ind in range(T):
        for om_ind in range(M_ome):
            scan = scanRange[scan_ind]
            frm = sf.wrapFrame(frmRange[om_ind])
            
            print(f'Scan {scan}, Frame {frm}')
            # 1. Load ROI
            roi = sf.loadROI(dataPath,scan,frm,etaRoi,tthRoi,params)
            roiTensor[:,:,om_ind,scan_ind] = roi
            
    return roiTensor

def saveROItensors(dome,output_path,spotInds,spotData,scanRange,dataFile,ttPath,spotsFiles,params):
    os.makedirs(output_path, exist_ok=True)
    for spotInd in spotInds:
        print(f'Spot {spotInd}')
        outFile = os.path.join(output_path,f'roiTensor_{spotInd}.pkl')
        roiTensor_k = sf.roiTensor(spotInd,spotData,dome,scanRange,dataFile,ttPath,ttPath,spotsFiles,params)
        with open(outFile, 'wb') as f:
            pickle.dump(roiTensor_k,f)
                
def makeTrackImages(dome,num_cols,output_path,spotInds,spotData,scanRange,dataFile,ttPath,spotsFiles,params):
    os.makedirs(output_path, exist_ok=True) 
        
    for spotInd in spotInds:
        fig_list = sf.roiTrackVisual(spotInd,spotData,dome,num_cols,scanRange,dataFile,ttPath,ttPath,spotsFiles,params)
        for i, fig in enumerate(fig_list):
            figure_file = os.path.join(output_path, f"fig_{spotInd}-{i}.png")
            fig.savefig(figure_file)          
            
def makeBlobTrackImages(num_cols,output_path,spotInds,spotData,scanRange,dataFileSequence,ttPath,spotsFiles,params):
    os.makedirs(output_path, exist_ok=True)        
    for spotInd in spotInds:
        fig_list = sf.showRoi3D(spotInd,spotData,num_cols,scanRange,dataFileSequence,ttPath,ttPath,spotsFiles,params)
        for i, fig in enumerate(fig_list):
            figure_file = os.path.join(output_path, f"fig_{spotInd}-{i}.png")
            fig.savefig(figure_file)  
        
        fig_list = sf.showRoiBlobs(spotInd,spotData,num_cols,scanRange,dataFileSequence,ttPath,ttPath,spotsFiles,params)
        for i, fig in enumerate(fig_list):
            figure_file = os.path.join(output_path,f'blobs_{spotInd}-{i}.png')
            fig.savefig(figure_file) 
            

def makeTrackSlides(ppt_file,output_path,spotInds,resultsData,spotData):
    ppt = Presentation()
    for si,spotInd in enumerate(spotInds):
        
        os.makedirs(output_path, exist_ok=True)
        # Save the current figure
        figure_file = os.path.join(output_path, f"fig_{spotInd}.png")
        
        # Slide 1 with track figure
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # Blank slide layout
        slide.shapes.add_picture(figure_file, Inches(0), Inches(0), Inches(10), Inches(8))
        
        
        # Slide 2 with track, data info 
        grNum = spotData['grain_nums'][spotInd]
        eta = spotData['etas'][spotInd]
        tth = spotData['tths'][spotInd]
        ome = spotData['omes'][spotInd]
        H = spotData['H'][spotInd]
        K = spotData['K'][spotInd]
        L = spotData['L'][spotInd]
        
        labels = ["truthDist", "changeDist", "omegaCount", "trackTimes"]
        table_vals = []
        for i in range(4):
            table_vals.append([])
            for t in range(5):
                table_vals[i].append(resultsData[labels[i]][si][t])
            
        slide2 = ppt.slides.add_slide(ppt.slide_layouts[6])  # Blank slide
        
        rows, cols = 2, 5  # Rows: labels + header, Columns: 6 (labels + 5 time steps)
        x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(1)  # Table position and size
        table = slide2.shapes.add_table(rows, cols, x, y, cx, cy).table
        table.cell(0,0).text = "grain"
        table.cell(0,1).text = "HKL"
        table.cell(0,2).text = "eta"
        table.cell(0,3).text = "tth"
        table.cell(0,4).text = "ome"
        table.cell(1,0).text = f"{grNum:.0f}"
        table.cell(1,1).text = f"{H:1.0f} {K:1.0f} {L:1.0f}"
        table.cell(1,2).text = f"{eta:1.4f}"
        table.cell(1,3).text = f"{tth:1.4f}"
        table.cell(1,4).text = f"{ome:1.4f}"
        
        rows, cols = len(labels) + 1, 6  # Rows: labels + header, Columns: 6 (labels + 5 time steps)
        x, y, cx, cy = Inches(1), Inches(3), Inches(8), Inches(4)  # Table position and size
        table = slide2.shapes.add_table(rows, cols, x, y, cx, cy).table
        
        # Populate header row
        for col in range(1, cols):
            table.cell(0, col).text = f"state {col - 1}"
        
        # Populate the data rows
        for row_idx, label in enumerate(labels, start=1):
            table.cell(row_idx, 0).text = label
            for col_idx, value in enumerate(table_vals[row_idx - 1], start=1):
                table.cell(row_idx, col_idx).text = format_row(row_idx,value)
        
    ppt.save(ppt_file)
    print(f"Figure added to PowerPoint: {ppt_file}")

def format_row(row_idx, value):
    if row_idx == 1:  # Row 1: Scientific Notation
        return f"{value:.2e}"
    elif row_idx == 2:  # Row 2: Sci
        return f"{value:.2e}"
    elif row_idx == 3:  # Row 3: Integer
        return f"{value:1.0f}"
    elif row_idx == 4:  # Row 4: Scientific notation
        return f"{value:.2f}"
    else:  # Default
        return str(value)

def loadSpotsAtFrame(spot_data,fnames,frame,params,detectFrame=[]):
    tths = spot_data['tths']
    etas = spot_data['etas']     
    ome_idxs = spot_data['ome_idxs'] 
    
    # Get spot indices at frame
    spotInds = np.where(ome_idxs == frame)[0]
    
    # Init ROIs list so they can be different sizes
    roi_list = []
    
    for i, ind in enumerate(spotInds):
       
        # 1. Load spot information
        tth = tths[ind]
        eta = etas[ind]

        # 2. Load polar ROI
        if detectFrame == []:
            detectFrame = frame
            
        roi_polar = sf.loadPolarROI([fnames],tth,eta,detectFrame,params)    
        
        roi_list.append(roi_polar)
        
    return roi_list



def plotROI(ax,roi,tthRoi,etaRoi,params):
        
    detectDist, mmPerPixel, ff_trans, ff_tilt = sf.loadYamlData(params,tthRoi,etaRoi)
    rad_dom, eta_dom = sf.polarDomain(detectDist, mmPerPixel,tthRoi, etaRoi, params['roiSize'])
    tth_dom = np.arctan(rad_dom*mmPerPixel/detectDist)
    
    roiPlt = ax.imshow(roi)
    plt.xticks([0,params['roiSize'][0]], [f'{eta_dom[0]:.4g}',f'{eta_dom[-1]:.4g}'])
    plt.yticks([0,params['roiSize'][1]], [f'{tth_dom[0]:.4g}',f'{tth_dom[-1]:.4g}'])
    ax.text(1, 2, f'max: {roi.max():.2f}', color='white', fontsize=12, weight='bold')
    
    return roiPlt

def showROI(ax,dataPath,scan,frm,tthRoi,etaRoi,params):
    
    roi = sf.loadROI(dataPath,scan,frm,etaRoi,tthRoi,params)
        
    detectDist, mmPerPixel, ff_trans, ff_tilt = sf.loadYamlData(params,tthRoi,etaRoi)
    rad_dom, eta_dom = sf.polarDomain(detectDist, mmPerPixel,tthRoi, etaRoi, params['roiSize'])
    tth_dom = np.arctan(rad_dom*mmPerPixel/detectDist)
    
    ax.imshow(roi)
    plt.xticks([0,params['roiSize'][0]], [f'{eta_dom[0]:.4g}',f'{eta_dom[-1]:.4g}'])
    plt.yticks([0,params['roiSize'][1]], [f'{tth_dom[0]:.4g}',f'{tth_dom[-1]:.4g}'])
  
    ax.text(1, 2, f'max: {roi.max():.2f}', color='white', fontsize=12, weight='bold')

def showROIpolar(ax,dataPath,scan,frm,tthRoi,etaRoi,params):
    
    roi = sf.loadROI(dataPath,scan,frm,etaRoi,tthRoi,params)
        
    detectDist, mmPerPixel, ff_trans, ff_tilt = sf.loadYamlData(params,tthRoi,etaRoi)
    rad_dom, eta_dom = sf.polarDomain(detectDist, mmPerPixel,\
                          tthRoi, etaRoi, params['roiSize'])
    tth_dom = np.arctan(rad_dom*mmPerPixel/detectDist)
    
    ax.imshow(roi)
    plt.xticks([0,39], [f'{eta_dom[0]:.4g}',f'{eta_dom[-1]:.4g}'])
    plt.yticks([0,39], [f'{tth_dom[0]:.4g}',f'{tth_dom[-1]:.4g}'])
  
    ax.text(1, 2, f'max: {roi.max():.2f}', color='white', fontsize=12, weight='bold')

def showInitial(ax,etaRoi,tthRoi,eta0,tth0,params):
    # Hexrd location
    [y_pos, x_pos] = sf.etaTthToPix(eta0,tth0,etaRoi,tthRoi,params)
    if (y_pos >= 0) and (y_pos < params['roiSize'][0]) and\
       (x_pos >= 0) and (x_pos < params['roiSize'][1]): 
        
        ax.plot(x_pos,y_pos,marker='s',markersize=10,\
                fillstyle='none',color='white')

def checkTrack(track_data,scan_ind,scan,frm):
    trackFound = False
    om_ind = 0
    if len(track_data) > scan_ind:
        for om_ind, omTrack in enumerate(track_data[scan_ind]):
            cond1 = omTrack['scan'] == scan_ind
            cond2 = omTrack['frm'] == frm
            if cond1 and cond2:
                trackFound = True
                break
    return trackFound, om_ind

def checkTrackBlob(track_data,scan_ind,scan,frm):
    trackFound = False
    om_ind = 0
    if len(track_data) > scan_ind:
        for om_ind, omTrack in enumerate(track_data[scan_ind]):
            cond1 = omTrack['scan'] == scan_ind
            cond2 = omTrack['frm'] == frm
            if cond1 and cond2:
                trackFound = True
                break
    return trackFound, om_ind

def checkTruth(truth_data,scan_ind,scan,frm):
    truthFound = False
    om_ind = 0
    if len(truth_data[scan_ind]) > scan_ind:
        for om_ind, omTruth in enumerate(truth_data[scan_ind]):
            if len(omTruth) > 0:
                cond1 = omTruth['scan'] == scan
                cond2 = omTruth['frm'] == frm
                if cond1 and cond2:
                    truthFound = True
                    break
    return truthFound, om_ind

def showTrack(ax,track,etaRoi,tthRoi,eta0,tth0,params):
    eta = track['eta']
    tth = track['tth']
    etaPrev = track['etaRoi']
    tthPrev = track['tthRoi']
    # Current track
    [y_pos, x_pos] = sf.etaTthToPix(eta,tth,etaRoi,tthRoi,params)
    if (y_pos >= 0) and (y_pos < params['roiSize'][0]) and\
       (x_pos >= 0) and (x_pos < params['roiSize'][1]):
        # ax.plot(x_pos,y_pos,marker='s',markersize=8,\
        #         fillstyle='none',color='red')
        ax.plot(x_pos,y_pos,marker='o',markersize=13,fillstyle='none',color='red')

def showTruth(ax,truth,etaRoi,tthRoi,params):
    eta1 = truth['eta1']
    tth1 = truth['tth1']
    eta2 = truth['eta2']
    tth2 = truth['tth2']
    [y1, x1] = sf.etaTthToPix(eta1,tth1,etaRoi,tthRoi,params)
    [y2, x2] = sf.etaTthToPix(eta2,tth2,etaRoi,tthRoi,params)
    
    rect = plt.Rectangle((x1, y1), x2-x1, y2-y1,
                          linewidth=1, edgecolor='g', facecolor='none')
    ax.add_patch(rect)

def showTensor(x):
    M3 = x.shape[2]
    figEdge = int(np.ceil(np.sqrt(M3)))
    fig, axes = plt.subplots(figEdge, figEdge, figsize=(figEdge*2, figEdge), sharex=True, sharey=True)
    ax = axes.ravel()
    
    for j in range(M3):
        img = x[:,:,j]
        ax[j].imshow(img)
        ax[j].set_axis_off()
    plt.tight_layout()
    plt.show()